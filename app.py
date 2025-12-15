import streamlit as st
import google.generativeai as genai
import pypdf
from pptx import Presentation
import docx2txt
import json
import re

# --- Configuração da Página ---
st.set_page_config(page_title="Gerador de Quizzes Universitário", page_icon="🎓", layout="centered")

st.title("🎓 Estuda com IA: Gerador de Quizzes")
st.write("Carrega os materiais da aula e personaliza o teu teste.")

# --- Barra Lateral para Configuração ---
with st.sidebar:
    st.header("⚙️ Configurações")
    
    # Campo de API Key
    api_key = st.text_input("Insere a tua API Key da Google", type="password")
    st.markdown("[Obter Chave Gratuita](https://aistudio.google.com/app/apikey)")
    
    st.divider() 
    
    # 1. Seletor de Modelo
    modelo_escolhido = st.selectbox(
        "Modelo da IA", 
        ["gemini-2.5-flash", "gemini-2.5-pro"],
        index=0
    )
    
    # 2. Nível de Dificuldade
    dificuldade = st.selectbox(
        "Nível de Dificuldade",
        ["Fácil (Memorização)", "Médio (Aplicação)", "Difícil (Análise Crítica)"],
        index=1
    )
    
    # 3. Tipos de Perguntas
    tipos_perguntas = st.multiselect(
        "Tipos de Perguntas",
        ["Múltipla Escolha", "Verdadeiro ou Falso", "Associação de Colunas"],
        default=["Múltipla Escolha", "Verdadeiro ou Falso"]
    )
    
    # 4. Quantidade de Perguntas
    qtd_perguntas = st.slider("Número de Perguntas", 3, 20, 5)

    # 5. Número de Alternativas
    num_alternativas = st.slider(
        "Opções (apenas para Múltipla Escolha)",
        3, 6, 4
    )

# --- Funções de Leitura de Ficheiros ---
def ler_pdf(file):
    pdf_reader = pypdf.PdfReader(file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text() or ""
    return text

def ler_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def ler_docx(file):
    return docx2txt.process(file)

# --- Função para extrair letra da resposta ---
def extrair_letra(texto):
    """Extrai a letra da resposta (A, B, C, etc.) de forma robusta"""
    if not texto:
        return None
    
    texto = str(texto).strip()
    
    # Se já for só uma letra
    if len(texto) == 1 and texto.isalpha():
        return texto.upper()
    
    # Se tiver formato "A)" ou "A) texto"
    match = re.match(r'^([A-Z])\)', texto, re.IGNORECASE)
    if match:
        return match.group(1).upper()
    
    # Se começar com letra seguida de qualquer coisa
    if texto[0].isalpha():
        return texto[0].upper()
    
    return None

# --- Função DEFINITIVA para processar SQL ---
def processar_pergunta_com_sql(pergunta_texto):
    """
    Processa texto e separa código SQL de forma inteligente.
    Retorna lista de tuplas: [('texto', conteudo), ('sql', codigo), ...]
    """
    
    # Substitui \n por quebras reais
    pergunta_texto = pergunta_texto.replace('\\n', '\n')
    
    # ESTRATÉGIA 1: Usa marcadores ```sql ... ```
    if '```sql' in pergunta_texto.lower() or '```' in pergunta_texto:
        partes = []
        # Split por ``` mas mantém o delimitador
        segmentos = re.split(r'(```(?:sql)?)', pergunta_texto, flags=re.IGNORECASE)
        
        dentro_codigo = False
        buffer_codigo = ""
        buffer_texto = ""
        
        for seg in segmentos:
            if re.match(r'```(?:sql)?', seg, re.IGNORECASE):
                if dentro_codigo:
                    # Fecha bloco de código
                    if buffer_codigo.strip():
                        partes.append(('sql', buffer_codigo.strip()))
                    buffer_codigo = ""
                    dentro_codigo = False
                else:
                    # Salva texto antes de abrir código
                    if buffer_texto.strip():
                        partes.append(('texto', buffer_texto.strip()))
                    buffer_texto = ""
                    # Abre bloco de código
                    dentro_codigo = True
            else:
                if dentro_codigo:
                    buffer_codigo += seg
                else:
                    buffer_texto += seg
        
        # Adiciona texto final se houver
        if buffer_texto.strip():
            partes.append(('texto', buffer_texto.strip()))
        
        return partes if partes else [('texto', pergunta_texto)]
    
    # ESTRATÉGIA 2: Detecta blocos SQL por keywords (CREATE, SELECT, INSERT completos)
    # Procura por comandos SQL completos terminados em ;
    sql_block_pattern = r'((?:CREATE\s+TABLE|SELECT|INSERT\s+INTO|UPDATE|DELETE\s+FROM)[^;]*;)'
    
    partes = []
    ultimo_fim = 0
    
    for match in re.finditer(sql_block_pattern, pergunta_texto, re.IGNORECASE | re.DOTALL):
        # Adiciona texto antes do SQL
        texto_antes = pergunta_texto[ultimo_fim:match.start()].strip()
        if texto_antes:
            partes.append(('texto', texto_antes))
        
        # Adiciona o bloco SQL
        sql_code = match.group(1).strip()
        partes.append(('sql', sql_code))
        
        ultimo_fim = match.end()
    
    # Adiciona texto restante
    texto_final = pergunta_texto[ultimo_fim:].strip()
    if texto_final:
        partes.append(('texto', texto_final))
    
    return partes if partes else [('texto', pergunta_texto)]

# --- Lógica Principal ---
st.subheader("1. Carregar Material")
uploaded_file = st.file_uploader("Arrasta o teu ficheiro aqui", type=['pdf', 'pptx', 'docx'])

tema_foco = st.text_input(
    "Queres focar num tema específico? (Opcional)",
    placeholder="Ex: Foca-te apenas nas datas históricas"
)

if uploaded_file is not None and api_key:
    # Extrair texto
    texto_extraido = ""
    try:
        if uploaded_file.name.endswith('.pdf'):
            texto_extraido = ler_pdf(uploaded_file)
        elif uploaded_file.name.endswith('.pptx'):
            texto_extraido = ler_pptx(uploaded_file)
        elif uploaded_file.name.endswith('.docx'):
            texto_extraido = ler_docx(uploaded_file)
        
        st.info(f"📄 Ficheiro carregado! ({len(texto_extraido)} caracteres)")
        
        if not tipos_perguntas:
            st.warning("⚠️ Por favor seleciona pelo menos um tipo de pergunta na barra lateral.")
        
        elif st.button("🚀 Gerar Quiz Personalizado", type="primary"):
            with st.spinner("A IA está a gerar as perguntas..."):
                
                genai.configure(api_key=api_key)
                model = genai.GenerativeModel(modelo_escolhido)

                # --- PROMPT ULTRA-ESPECÍFICO ---
                prompt = f"""
Atua como um professor universitário experiente. Cria EXATAMENTE {qtd_perguntas} perguntas de quiz baseadas neste conteúdo:

CONTEÚDO DO MATERIAL:
"{texto_extraido[:30000]}"

⚠️ CONFIGURAÇÕES OBRIGATÓRIAS:
- Quantidade: EXATAMENTE {qtd_perguntas} perguntas (nem mais, nem menos)
- Dificuldade: {dificuldade}
- Foco específico: {tema_foco if tema_foco else "Todos os tópicos do material"}
- Tipos de perguntas permitidos: {', '.join(tipos_perguntas)}
- Número de alternativas (múltipla escolha): {num_alternativas}

🔴 REGRA CRÍTICA DE FORMATAÇÃO SQL:
Quando incluíres código SQL, tabelas ou dados na pergunta, usa OBRIGATORIAMENTE este formato:

EXEMPLO CORRETO:
"Considere as seguintes tabelas:\\n\\n```sql\\nCREATE TABLE Equipas (\\n    idEquipa INT PRIMARY KEY,\\n    nome VARCHAR(100)\\n);\\n```\\n\\nDados inseridos:\\n\\n```sql\\nINSERT INTO Equipas (nome) VALUES ('Porto'), ('Benfica');\\n```\\n\\nQual o resultado da query:\\n\\n```sql\\nSELECT * FROM Equipas;\\n```"

📋 REGRAS DE FORMATAÇÃO POR TIPO:

1. **Múltipla Escolha**:
   - {num_alternativas} opções no formato: "A) texto", "B) texto", etc.
   - resposta_correta: APENAS a letra (ex: "A")
   
2. **Verdadeiro/Falso**:
   - Opções: ["A) Verdadeiro", "B) Falso"]
   - resposta_correta: "A" ou "B"

3. **Associação de Colunas**:
   - Formato: "Associe os itens:\\n\\n1. Item A\\n2. Item B\\n\\n--- Separador ---\\n\\nA. Definição X\\nB. Definição Y"
   - Opções com combinações: ["A) 1-A, 2-B", "B) 1-B, 2-A", ...]
   - resposta_correta: letra da combinação correta

📊 IMPORTANTE SOBRE CONTEXTO:
- Cada pergunta deve ser AUTOCONTIDA (incluir TODOS os dados necessários)
- Se a pergunta precisa de tabelas, dados ou código, INCLUI TUDO no campo 'pergunta'
- O aluno NÃO tem acesso ao material original durante o teste
- Usa \\n para quebras de linha dentro das strings JSON
- TODO código SQL deve estar entre ```sql e ```

✅ FORMATO JSON OBRIGATÓRIO (retorna APENAS isto, sem texto adicional):
[
    {{
        "tipo": "Múltipla Escolha",
        "pergunta": "Texto introdutório.\\n\\n```sql\\nCREATE TABLE exemplo (id INT);\\n```\\n\\nQual a função?",
        "opcoes": ["A) opção1", "B) opção2", "C) opção3", "D) opção4"],
        "resposta_correta": "A",
        "explicacao": "Explicação detalhada da resposta correta"
    }},
    {{
        "tipo": "Verdadeiro ou Falso",
        "pergunta": "O comando DROP apaga tabelas permanentemente.",
        "opcoes": ["A) Verdadeiro", "B) Falso"],
        "resposta_correta": "A",
        "explicacao": "DROP remove a tabela e todos os dados de forma irreversível."
    }}
]

🔍 VALIDAÇÃO FINAL ANTES DE RESPONDER:
1. Conta as perguntas: devem ser EXATAMENTE {qtd_perguntas}
2. Verifica se cada 'resposta_correta' é uma letra simples (A, B, C, D...)
3. Verifica se todo código SQL está entre ```sql e ```
4. Verifica se cada pergunta inclui TODOS os dados necessários
5. Verifica se o JSON é válido (sem vírgulas extras, aspas corretas)

⚠️ LEMBRA-TE: Retorna um array JSON com EXATAMENTE {qtd_perguntas} objetos!
"""
                
                try:
                    response = model.generate_content(
                        prompt,
                        generation_config={
                            "response_mime_type": "application/json",
                            "temperature": 0.7,
                        }
                    )
                    
                    texto_resposta = response.text.replace("```json", "").replace("```", "").strip()
                    
                    inicio = texto_resposta.find('[')
                    fim = texto_resposta.rfind(']') + 1

                    if inicio != -1 and fim != 0:
                        json_str = texto_resposta[inicio:fim]
                        quiz_data = json.loads(json_str)
                        
                        # ✅ VALIDAÇÃO E CORREÇÃO DO NÚMERO DE PERGUNTAS
                        if len(quiz_data) > qtd_perguntas:
                            quiz_data = quiz_data[:qtd_perguntas]
                            st.warning(f"⚠️ A IA gerou mais perguntas. Foram cortadas para {qtd_perguntas}.")
                        elif len(quiz_data) < qtd_perguntas:
                            st.warning(f"⚠️ A IA gerou apenas {len(quiz_data)} perguntas (pediste {qtd_perguntas}).")
                        
                        # Validação e limpeza dos dados
                        quiz_limpo = []
                        for idx, q in enumerate(quiz_data):
                            # Garante que todos os campos existem
                            if all(key in q for key in ['tipo', 'pergunta', 'opcoes', 'resposta_correta', 'explicacao']):
                                # Limpa a resposta_correta
                                q['resposta_correta'] = extrair_letra(q['resposta_correta']) or "A"
                                
                                # Garante que opcoes é uma lista
                                if not isinstance(q['opcoes'], list):
                                    st.warning(f"⚠️ Pergunta {idx+1} tem opções inválidas. Ignorada.")
                                    continue
                                
                                quiz_limpo.append(q)
                            else:
                                st.warning(f"⚠️ Pergunta {idx+1} está incompleta. Ignorada.")
                        
                        if quiz_limpo:
                            st.session_state['quiz_data'] = quiz_limpo
                            
                            # Limpar estados antigos
                            for key in list(st.session_state.keys()):
                                if key.startswith('q_') or key.startswith('respondido_'):
                                    del st.session_state[key]
                            
                            st.success(f"✅ Quiz gerado com {len(quiz_limpo)} perguntas!")
                            st.rerun()
                        else:
                            st.error("❌ Nenhuma pergunta válida foi gerada. Tenta novamente.")
                    else:
                        st.error("❌ A IA não devolveu JSON válido. Tenta novamente.")

                except json.JSONDecodeError as e:
                    st.error(f"❌ Erro ao processar JSON: {e}")
                    with st.expander("🔍 Ver resposta da IA (debug)"):
                        st.code(texto_resposta)
                except Exception as e:
                    st.error(f"❌ Erro na API: {e}")

    except Exception as e:
        st.error(f"❌ Erro ao ler ficheiro: {e}")

# --- MOSTRAR O QUIZ (RENDERIZAÇÃO CORRIGIDA) ---
if 'quiz_data' in st.session_state:
    st.markdown("---")
    st.subheader(f"📝 Quiz Gerado ({len(st.session_state['quiz_data'])} Perguntas)")
    
    respostas_certas = 0
    respostas_dadas = 0
    total = len(st.session_state['quiz_data'])
    
    for i, q in enumerate(st.session_state['quiz_data']):
        tipo_label = q.get('tipo', 'Pergunta')
        
        # Container para cada pergunta
        with st.container():
            st.markdown(f"### 📌 Pergunta {i+1} de {total}")
            st.caption(f"**Tipo:** {tipo_label}")
            
            # --- RENDERIZAÇÃO INTELIGENTE ---
            texto_pergunta = q['pergunta']
            
            # CASO ESPECIAL: Associação de colunas
            if "--- Separador ---" in texto_pergunta:
                partes = texto_pergunta.split("--- Separador ---")
                col1, col2 = st.columns(2)
                
                with col1:
                    st.markdown("**Coluna 1:**")
                    for tipo, conteudo in processar_pergunta_com_sql(partes[0]):
                        if tipo == 'sql':
                            st.code(conteudo, language='sql')
                        else:
                            st.markdown(conteudo)
                
                with col2:
                    st.markdown("**Coluna 2:**")
                    for tipo, conteudo in processar_pergunta_com_sql(partes[1]):
                        if tipo == 'sql':
                            st.code(conteudo, language='sql')
                        else:
                            st.markdown(conteudo)
            
            # CASO NORMAL: Pergunta com ou sem SQL
            else:
                partes = processar_pergunta_com_sql(texto_pergunta)
                
                for tipo, conteudo in partes:
                    if tipo == 'sql':
                        st.code(conteudo, language='sql')
                    else:
                        st.markdown(conteudo)
            
            # Opções de resposta
            escolha = st.radio(
                "**Seleciona a tua resposta:**", 
                q['opcoes'], 
                key=f"q_{i}", 
                index=None
            )
            
            # Verificação da resposta
            if escolha:
                # Marca como respondida
                if f'respondido_{i}' not in st.session_state:
                    st.session_state[f'respondido_{i}'] = True
                    respostas_dadas += 1
                
                letra_user = extrair_letra(escolha)
                letra_correta = extrair_letra(q.get('resposta_correta', ''))
                
                if letra_user and letra_correta:
                    if letra_user == letra_correta:
                        st.success(f"✅ **Correto!**")
                        st.info(f"💡 **Explicação:** {q.get('explicacao', 'Sem explicação disponível.')}")
                        if f'certa_{i}' not in st.session_state:
                            st.session_state[f'certa_{i}'] = True
                            respostas_certas += 1
                    else:
                        st.error(f"❌ **Errado.** A resposta correta era: **{letra_correta})**")
                        st.info(f"💡 **Explicação:** {q.get('explicacao', 'Sem explicação disponível.')}")
                else:
                    st.warning("⚠️ Erro ao processar a resposta. Por favor reporta este bug.")
            
            st.markdown("---")

    # Contar respostas certas do session_state
    respostas_certas = sum(1 for key in st.session_state.keys() if key.startswith('certa_'))
    respostas_dadas = sum(1 for key in st.session_state.keys() if key.startswith('respondido_'))

    # Resultado final
    if total > 0:
        percentagem = (respostas_certas / total) * 100 if respostas_dadas > 0 else 0
        
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("✅ Respostas Certas", f"{respostas_certas}/{total}")
        with col2:
            st.metric("📊 Percentagem", f"{percentagem:.0f}%")
        with col3:
            st.metric("📝 Respondidas", f"{respostas_dadas}/{total}")
        
        if respostas_dadas == total:
            if respostas_certas == total:
                st.balloons()
                st.success("🎉 **PERFEITO! Acertaste todas as perguntas!**")
            elif percentagem >= 70:
                st.success("👏 **Bom trabalho! Passaste no teste!**")
            elif percentagem >= 50:
                st.info("📚 **Razoável. Revê alguns tópicos e tenta novamente.**")
            else:
                st.warning("💪 **Continua a estudar! Não desistas, vais conseguir!**")

elif not api_key:
    st.warning("👈 Insere a tua API Key da Google na barra lateral para começar.")
else:
    st.info("📤 Carrega um ficheiro (PDF, PPTX ou DOCX) para gerar o teu quiz personalizado.")
